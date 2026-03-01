import json
import os
import re
import shutil
import subprocess
import tempfile
import csv
from datetime import datetime
from typing import List, Optional, Dict, Any

import PyPDF2
import docx
import pdfplumber
from pptx import Presentation

from models.chunk import TextChunk
import openpyxl

from utils.logger import logger
import pandas as pd
from config.config import BASE_DIR


class MarkdownNode:
    def __init__(self, level, content, parent=None):
        self.level = level  # 1-6 for headers, 99 for body/table, 98 for bound header+table
        self.content = content
        self.children = []
        self.parent = parent
        self.metadata = {}  # 存储额外信息，如 'header_text' 用于长表格切分时的表头复制
        
    @property
    def length(self):
        # 自身长度 + 所有子节点长度
        return len(self.content) + sum(c.length for c in self.children)
    
    def get_full_text(self):
        # 递归获取完整文本
        parts = [self.content]
        for child in self.children:
            parts.append(child.get_full_text())
        return "\n".join(parts)


class MarkdownSemanticSplitter:
    def __init__(self, chunk_size, overlap, language, file_name, scan_items):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.language = language
        self.file_name = file_name
        self.scan_items = scan_items
        # 噪声过滤模式
        self.noise_patterns = [
            r'^\s*\d+\s*$',  # 纯数字页码
            r'^文档版本\s+\d+',  # 版本号
            r'^发布日期\s+\d{4}-\d{2}-\d{2}',  # 日期
            r'^版权所有\s+©',  # 版权
            r'^非经本公司书面许可',
            r'^华为云计算技术有限公司',
        ]
        # 表头匹配模式：匹配 "表 1-1" 或 "Table 1" 等
        self.table_caption_pattern = r'^(表|Table)\s*\d+[-.]?\d*.*'

    def split_markdown_file(self, md_path) -> List[TextChunk]:
        # Legacy method kept for compatibility if needed, redirect to v2
        return self.split_markdown_file_v2(md_path)

    def _parse_to_tree(self, lines: List[str]) -> MarkdownNode:
        root = MarkdownNode(level=0, content="")
        current_node = root
        # 维护一个路径栈：[root, level1, level2...]
        stack = [root]

        in_table = False
        table_lines = []
        
        # 暂存上一个段落节点，用于检测 "表头 + 表格" 的绑定关系
        last_body_node = None 
        
        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 1. 噪声过滤
            if any(re.search(p, line) for p in self.noise_patterns):
                continue

            # 2. 表格处理 (HTML Table)
            if '<table>' in line:
                in_table = True
                table_lines.append(line)
                # 修复：检查单行完整表格（<table>...</table>在同一行）
                if '</table>' in line:
                    in_table = False
                    table_content = "\n".join(table_lines)
                    
                    # 检测表头绑定
                    table_caption = ""
                    if last_body_node and last_body_node.level == 99:
                        if re.match(self.table_caption_pattern, last_body_node.content):
                            table_caption = last_body_node.content
                            if last_body_node.parent:
                                last_body_node.parent.children.remove(last_body_node)
                            combined_content = f"{table_caption}\n{table_content}"
                            table_node = MarkdownNode(level=99, content=combined_content, parent=current_node)
                            table_node.metadata['header_text'] = table_caption
                            current_node.children.append(table_node)
                            last_body_node = table_node
                            table_lines = []
                            continue
                    
                    # 普通单行表格
                    table_node = MarkdownNode(level=99, content=table_content, parent=current_node)
                    current_node.children.append(table_node)
                    last_body_node = table_node
                    table_lines = []
                continue
            if in_table:
                table_lines.append(line)
                if '</table>' in line:
                    in_table = False
                    table_content = "\n".join(table_lines)
                    
                    # 关键逻辑优化：检测表头绑定
                    # 如果上一个节点是 body_node (level=99)，且内容匹配表头模式
                    table_caption = ""
                    if last_body_node and last_body_node.level == 99:
                         if re.match(self.table_caption_pattern, last_body_node.content):
                            # 命中！将上一个节点标记为 "待合并表头"
                            # 但实际上，更简单的做法是：
                            # 1. 从父节点中移除 last_body_node
                            # 2. 创建一个新的 "绑定节点" (level=98)，包含表头
                            # 3. 将表格作为绑定节点的子节点或直接合并内容？
                            # 为了保持树结构清晰，我们不仅合并内容，而是将它们视为一个逻辑单元。
                            
                            # 策略：创建一个组合内容节点
                            table_caption = last_body_node.content
                            
                            # 从父节点的 children 中移除 last_body_node
                            if last_body_node.parent:
                                last_body_node.parent.children.remove(last_body_node)
                            
                            # 将表头内容拼接到表格前面，作为一个整体节点
                            # 或者，我们仍然保持它们分开，但在 _recursive_split 时特殊处理？
                            # 最稳妥的是物理合并：Text = Caption + \n + Table
                            # 注意：我们不再将 caption 塞入 Table 标签内，而是作为文本前置
                            combined_content = f"{table_caption}\n{table_content}"
                            table_node = MarkdownNode(level=99, content=combined_content, parent=current_node)
                            # 存储表头信息，供后续超长切分时使用
                            table_node.metadata['header_text'] = table_caption
                            current_node.children.append(table_node)
                            last_body_node = table_node # 更新 last_body_node
                            table_lines = []
                            continue

                    # 如果没有命中表头，则作为普通表格节点添加
                    table_node = MarkdownNode(level=99, content=table_content, parent=current_node)
                    current_node.children.append(table_node)
                    last_body_node = table_node
                    table_lines = []
                continue

            # 3. 标题处理 (# Header)
            header_match = re.match(r'^(#+)\s+(.*)', line)
            if header_match:
                level = len(header_match.group(1))
                content = header_match.group(2).strip()
                
                # 找到正确的父节点
                # 如果当前level比栈顶level深，说明是子节点，直接入栈
                # 如果当前level比栈顶浅或相等，说明前面的子树结束了，出栈直到找到父节点
                while len(stack) > 1 and stack[-1].level >= level:
                    stack.pop()
                
                parent = stack[-1]
                new_node = MarkdownNode(level=level, content=line, parent=parent) # 保留完整标题行 # Header
                parent.children.append(new_node)
                stack.append(new_node)
                current_node = new_node
                last_body_node = None # 标题打断了 body 序列
            else:
                # 4. 普通段落 / 图片
                # 归属于当前栈顶节点
                body_node = MarkdownNode(level=99, content=line, parent=stack[-1])
                stack[-1].children.append(body_node)
                last_body_node = body_node

        return root

    def _get_len(self, acc):
        return sum(len(s) for s in acc) + len(acc) # + len(acc) for newlines roughly

    def _flush_chunk(self, acc, chunks):
        if not acc:
            return
        text = "\n".join(acc)
        chunks.append(TextChunk(
            file_path=self.file_name, # 注意这里file_name传了full path
            text=text,
            length=len(text),
            language=self.language,
            file_name=os.path.basename(self.file_name),
            scan_items=self.scan_items,
            metadata={"strategy": "adaptive_markdown"}
        ))
        acc.clear()

    def _split_and_flush_long_text(self, text, chunks, header_text=""):
        """
        处理超长文本：尝试按句子安全切分，如果单句仍然超长，再退回到硬切分。
        逻辑复用 DocumentSplitter.split_long_paragraph 的思想。
        增加 header_text 参数：如果是切割长表格，需要把表头带入每个子块。
        """
        # 复用 split_long_paragraph 的逻辑
        if len(text) <= self.chunk_size:
            final_text = text
            # 注意：移除之前可能存在的对 header_text 的不当处理
            
            chunks.append(TextChunk(
                file_path=self.file_name,
                text=final_text,
                length=len(final_text),
                language=self.language,
                file_name=os.path.basename(self.file_name),
                scan_items=self.scan_items,
                metadata={"strategy": "adaptive_markdown"}
            ))
            return

        # 检测是否是表格
        is_table = '<table>' in text
        
        # 尝试按句子分割
        sentences = []
        
        if is_table:
            # 如果是表格，按行 </tr> 分割
            # 这是一个简单的 HTML 表格分割器
            # 提取所有 tr
            # 1. 提取 table 标签属性（如果有）
            table_start_match = re.search(r'<table[^>]*>', text)
            table_start_tag = table_start_match.group(0) if table_start_match else "<table>"
            
            # 使用 regex 提取所有行
            tr_pattern = re.compile(r'<tr.*?>.*?</tr>', re.DOTALL)
            all_rows = tr_pattern.findall(text)
            
            header_row = ""
            preamble = ""
            
            if not all_rows:
                # 居然没匹配到 tr，退化为普通文本处理
                is_table = False
            else:
                sentences = all_rows
                header_row = all_rows[0]
                
                # 提取 preamble (表格标签及其前面的所有内容，包括标题)
                # 使用 finditer 找到第一个 tr 的位置
                first_match = next(tr_pattern.finditer(text))
                preamble = text[:first_match.start()]
        
        if not is_table:
            current_sentence = ""
            for char in text:
                current_sentence += char
                if char in ('。', '？', '！', '\n'):
                    sentences.append(current_sentence)
                    current_sentence = ""
            if current_sentence:
                sentences.append(current_sentence)
        
        current_chunk = ""
        # 如果是表格，当前 chunk 初始化需要包含 table start tag 和 header row (对于非第一块)
        # 但第一块自然会包含。
        # 我们用一个 flag 指示是否是第一块
        is_first_chunk = True
        
        if is_table:
            current_chunk = preamble
        
        # 预计算表格外壳开销
        table_overhead = 0
        if is_table:
             # <table> + header_row + </table> + header_text(续表)
             overhead_str = f"{header_text} (续表)\n{table_start_tag}{header_row}</table>"
             table_overhead = len(overhead_str)
        
        for i, sentence in enumerate(sentences):
            # sentence 此时是单句或单行 <tr>...</tr>
            
            # 如果是表格模式，sentence 就是一行 tr
            token_len = len(sentence)
            
            # 检查加入后是否超长
            # 如果是表格模式，我们要预留闭合标签的空间 </table>
            limit = self.chunk_size
            if is_table:
                limit -= 8 # len('</table>')
            
            # [Fix] 强制兜底：如果单行/单句本身就超过 limit，必须强制切分
            if token_len > limit:
                # 1. Flush current buffer if any
                if current_chunk:
                    final_text = current_chunk
                    if is_table: final_text += "</table>"
                    chunks.append(TextChunk(
                        file_path=self.file_name,
                        text=final_text,
                        length=len(final_text),
                        language=self.language,
                        file_name=os.path.basename(self.file_name),
                        scan_items=self.scan_items,
                        metadata={"strategy": "adaptive_markdown"}
                    ))
                    current_chunk = ""
                    is_first_chunk = False
                
                # 2. Hard split the long sentence
                # Note: This might break HTML structure for table rows, but length constraint is hard.
                start = 0
                while start < token_len:
                    end = min(start + limit, token_len)
                    sub_text = sentence[start:end]
                    
                    chunks.append(TextChunk(
                        file_path=self.file_name,
                        text=sub_text,
                        length=len(sub_text),
                        language=self.language,
                        file_name=os.path.basename(self.file_name),
                        scan_items=self.scan_items,
                        metadata={"strategy": "adaptive_markdown", "note": "hard_split"}
                    ))
                    start += limit
                
                # 3. Reset for next sentence (Restore table context if needed)
                if is_table:
                     prefix = f"{table_start_tag}{header_row}"
                     if header_text:
                         if "(续表)" not in header_text:
                            prefix = f"{header_text} (续表)\n{prefix}"
                         else:
                            prefix = f"{header_text}\n{prefix}"
                     current_chunk = prefix
                continue

            if len(current_chunk) + token_len > limit:
                # 当前块满了，需要 flush
                if current_chunk:
                    # 封装 chunk
                    final_text = current_chunk
                    if is_table:
                        # 补全表格结尾
                        final_text += "</table>"
                    
                    chunks.append(TextChunk(
                        file_path=self.file_name,
                        text=final_text,
                        length=len(final_text),
                        language=self.language,
                        file_name=os.path.basename(self.file_name),
                        scan_items=self.scan_items,
                        metadata={"strategy": "adaptive_markdown"}
                    ))
                    
                    # 开启新的一块
                    current_chunk = ""
                    is_first_chunk = False
                    
                    if is_table:
                        # 新块初始化：Header Text (续表) + Table Start + Header Row
                        prefix = f"{table_start_tag}{header_row}"
                        if header_text:
                             # 只有当 header_text 不包含 "(续表)" 时才添加，避免重复
                             if "(续表)" not in header_text:
                                prefix = f"{header_text} (续表)\n{prefix}"
                             else:
                                prefix = f"{header_text}\n{prefix}"
                        current_chunk = prefix

            # 添加当前 sentence
            current_chunk += sentence
        
        # 最后的残留
        if current_chunk:
            final_text = current_chunk
            if is_table:
                final_text += "</table>"

            chunks.append(TextChunk(
                file_path=self.file_name,
                text=final_text,
                length=len(final_text),
                language=self.language,
                file_name=os.path.basename(self.file_name),
                scan_items=self.scan_items,
                metadata={"strategy": "adaptive_markdown"}
            ))

    def _recursive_split(self, node: MarkdownNode, accumulator: List[str], chunks: List[TextChunk]):
        """
        accumulator: 当前正在构建的chunk内容列表
        chunks: 结果列表
        """
        # 定义短 chunk 阈值：如果 accumulator 太短，尽量避免单独 flush
        MIN_CHUNK_SIZE = self.chunk_size * 0.1
        
        # 1. 尝试将节点自身内容加入
        # 如果是根节点，跳过内容
        if node.level == 0:
            pass
        else:
            acc_len = self._get_len(accumulator)
            node_len = len(node.content)
            
            # 检查加入当前节点自身内容（标题或段落）后是否超长
            if acc_len + node_len > self.chunk_size:
                # 如果当前节点内容超长，需要特殊处理
                if node_len > self.chunk_size:
                    # 策略：如果 accumulator 较短，将其作为 prefix 与超长内容合并
                    # 避免产生孤立的短 chunk
                    header_text = node.metadata.get('header_text', "")
                    
                    if accumulator and acc_len < MIN_CHUNK_SIZE:
                        # accumulator 太短，作为前缀与超长内容合并
                        prefix = "\n".join(accumulator)
                        accumulator.clear()
                        # 将 prefix 拼接到 node.content 前面
                        combined_content = f"{prefix}\n{node.content}"
                        self._split_and_flush_long_text(combined_content, chunks, header_text=header_text)
                    else:
                        # accumulator 够长，正常 flush
                        if accumulator:
                            self._flush_chunk(accumulator, chunks)
                        self._split_and_flush_long_text(node.content, chunks, header_text=header_text)
                else:
                    # 当前节点内容不超长，但 accumulator + node 超长
                    # 策略：检查是否能"借位"——如果 accumulator 很短，允许轻微超长
                    if acc_len < MIN_CHUNK_SIZE:
                        # accumulator 太短，允许轻微超长以避免产生短 chunk
                        accumulator.append(node.content)
                    else:
                        # accumulator 够长，正常 flush 后添加 node
                        self._flush_chunk(accumulator, chunks)
                        accumulator.append(node.content)
            else:
                accumulator.append(node.content)

        # 2. 处理子节点
        # 贪心策略：直接递归，让递归过程中的 _get_len 判断来决定是否 flush。
        # 这样能自然地处理 "父标题 + 部分子节点" 在一个 chunk，"剩余子节点" 在下一个 chunk 的情况。
        
        for child in node.children:
            self._recursive_split(child, accumulator, chunks)

    # 替换 split_markdown_file 的实现
    def split_markdown_file_v2(self, md_path) -> List[TextChunk]:
        with open(md_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        root = self._parse_to_tree(lines)
        chunks = []
        accumulator = []
        self._recursive_split(root, accumulator, chunks)
        
        # 剩下的 flush
        if accumulator:
            self._flush_chunk(accumulator, chunks)
            
        return chunks


class DocumentSplitter():
    """文档拆分工具类"""

    def __init__(self, file_path: str="", chunk_size: int=1000, overlap: int=100,
                language="cn", file_name="", scan_items=None, adaptive_chunk=False, markdown_file_path: str = "", task_id=None):
        """
        初始化拆分器

        :param chunk_size: 每段的最大字符数
        :param overlap: 段落间的重叠字符数
        :param markdown_file_path: 指定预处理好的Markdown文件路径（用于自适应分块）
        :param task_id: 任务ID，用于生成日志文件路径
        """
        if scan_items is None:
            scan_items = []
        self.file_path = file_path
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.sentence_endings = r'(?<=[.!?。！？])\s+'
        self.language = language
        self.file_name = file_name
        self.scan_items = scan_items
        self.words = 0
        self.paragraphs = 0
        self.adaptive_chunk = adaptive_chunk
        self.markdown_file_path = markdown_file_path
        self.task_id = task_id

    def _log_chunk_stats(self, chunks: List[TextChunk], strategy: str, file_path: str):
        """
        记录分块统计信息到 CSV 文件
        """
        try:
            if self.task_id:
                 # 优先写入到 temp_tasks/{task_id}/chunk_stats.csv
                 from config.config import TEMP_DIR
                 stats_file = os.path.join(TEMP_DIR, self.task_id, "chunk_stats.csv")
            else:
                 # 回退到 logs/chunk_stats.csv
                 stats_file = os.path.join(str(BASE_DIR), "logs", "chunk_stats.csv")
            
            os.makedirs(os.path.dirname(stats_file), exist_ok=True)
            
            file_exists = os.path.exists(stats_file)
            
            with open(stats_file, 'a', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                if not file_exists:
                    writer.writerow(['chunk_index', 'file_name', 'strategy', 'chunk_length', 'total_chunks'])
                
                base_name = os.path.basename(file_path)
                
                for i, chunk in enumerate(chunks):
                    writer.writerow([i, base_name, strategy, chunk.length, len(chunks)])
                    
            logger.info(f"分块统计已写入: {stats_file}")
        except Exception as e:
            logger.error(f"写入分块统计失败: {e}")

    def split_document(self, file_path: str) -> List[TextChunk]:
        """
        根据文件类型自动选择处理方法

        :param file_path: 文件路径
        :return: 拆分后的文本段列表
        """
        file_path = os.path.abspath(file_path)
        ext = os.path.splitext(file_path)[1].lower()

        # 优先处理 Markdown 自适应分块逻辑
        if self.adaptive_chunk:
            # 获取文件名（不带后缀）
            base_name = os.path.splitext(os.path.basename(file_path))[0]
            
            # 直接去 config 中定义的 dirty_md 目录查找
            # 目录: BASE_DIR/files/dirty_md
            dirty_md_dir = os.path.join(str(BASE_DIR), "files", "dirty_md")
            
            # 1. 优先查找同名的 .md 文件
            potential_md_path = os.path.join(dirty_md_dir, f"{base_name}.md")
            
            # 2. 如果没找到，尝试查找 _dirty.md (兼容现有文件命名)
            if not os.path.exists(potential_md_path):
                potential_md_path = os.path.join(dirty_md_dir, f"{base_name}_dirty.md")
            
            # 3. 如果还是没找到，且用户显式传入了 markdown_file_path，则使用它
            if not os.path.exists(potential_md_path) and self.markdown_file_path:
                potential_md_path = self.markdown_file_path

            # 4. 如果本地都没有，尝试通过 MinerU API 在线解析
            if not os.path.exists(potential_md_path):
                logger.info(f"【自适应分块】本地未找到 Markdown 文件，尝试通过 MinerU API 在线解析: {base_name}")
                try:
                    from utils.mineru_api_util import parse_pdf_via_mineru_api
                    # 使用任务目录存放解析结果
                    if self.task_id:
                        mineru_output_dir = os.path.join(str(TEMP_DIR), self.task_id, "mineru_output")
                    else:
                        mineru_output_dir = os.path.join(str(BASE_DIR), "files", "mineru_output")
                    
                    lang = "cn" if self.language == "cn" else "en"
                    api_result = parse_pdf_via_mineru_api(file_path, mineru_output_dir, language=lang)
                    if api_result and os.path.exists(api_result):
                        potential_md_path = api_result
                        logger.info(f"【自适应分块】MinerU API 解析成功: {potential_md_path}")
                except Exception as e:
                    logger.warning(f"【自适应分块】MinerU API 解析失败: {e}")

            if os.path.exists(potential_md_path):
                logger.info(f"【自适应分块】传入 {os.path.basename(file_path)}，找到 {os.path.basename(potential_md_path)}，启用动态分块策略")
                splitter = MarkdownSemanticSplitter(
                    chunk_size=self.chunk_size,
                    overlap=self.overlap,
                    language=self.language,
                    file_name=file_path, # Keep original file path for reference
                    scan_items=self.scan_items
                )
                chunks = splitter.split_markdown_file_v2(potential_md_path)
                self._log_chunk_stats(chunks, "adaptive_markdown", file_path)
                return chunks
            else:
                logger.warning(f"【自适应分块】在 {dirty_md_dir} 下未找到对应的 Markdown 文件 ({base_name}.md 或 {base_name}_dirty.md)，且 MinerU API 解析失败，回退到固定分块 Baseline 策略。")

        chunks = []
        if ext in ['.xlsx', '.xls']:
            chunks = self._split_excel(file_path)
        elif ext == '.chm':
            chunks = self._split_chm(file_path)
        elif ext == '.pdf':
            chunks = self._split_pdf(file_path)
        elif ext == '.docx':
            chunks = self._split_docx(file_path)
        elif ext == '.pptx':
            chunks = self._split_pptx(file_path)
        elif ext == '.md': # Handle standard markdown splitting if adaptive is off
             with open(file_path, 'r', encoding='utf-8') as f:
                 chunks = self.split_text_into_chunks(f.read())
        else:
            raise ValueError(f"不支持的文件类型: {ext}")
            
        self._log_chunk_stats(chunks, "ours_fixed", file_path)
        return chunks
    
    def _split_pptx(self, file_path: str) -> List[TextChunk]:
        """处理pptx文件，提取文本和表格，表格转换为JSON格式"""
        try:
            # 读取pptx文件
            prs = Presentation(file_path)
            full_text = []

            # 遍历所有幻灯片
            for slide_index, slide in enumerate(prs.slides, 1):
                # 遍历幻灯片中的每个形状
                for shape in slide.shapes:
                    # 处理文本框内容
                    if shape.has_text_frame:
                        text_content = []
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                text_content.append(paragraph.text.strip())
                        
                        if text_content:
                            full_text.extend(text_content)

                    elif shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            # 提取单元格文本并清理
                            row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            
                            if row_data: # 跳过空行
                                # 将行数据转换为字典(第n列：值) - 与docx处理逻辑一致
                                row_dict = {}
                                for j, value in enumerate(row_data):
                                    row_dict[f"column {j+1}"] = value
                                table_data.append(row_dict)
                        
                        if table_data:
                            # 与docx相同的JSON格式
                            full_text.append(json.dumps(table_data, ensure_ascii=False))
            
            # 统计信息
            self.paragraphs = len(prs.slides)
            text_content = '\n'.join(full_text)
            self.words = len(text_content)

            # 调用分块方法
            return self.split_text_into_chunks(text_content)

        except Exception as e:
            print(f"处理pptx文件时发生错误: {str(e)}")
            return []
    
    def _split_pdf(self, file_path: str) -> List[TextChunk]:
        """处理pdf文件"""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            self.paragraphs = len(pdf_reader.pages)
        processor = PrecisionPDFProcessor()
        result = processor.process_pdf(
            file_path,
            text_clean_options={
                'remove_special_chars': True,
                'remove_extra_whitespace': False,
                'min_line_length': 3,
                'precise_newline_handling': False
            },
            table_clean_options={
                'header_rows': 1,
                'skip_footer': 0,
                'rename_columns': True,
                'precise_newline_handling': True
            }
        )
        text = '\n'.join(result)
        self.words = len(text)

        # 如果 self.adaptive_chunk 为 True 但代码执行到这里，
        # 说明没有找到对应的 Markdown 文件，回退到固定分块 (Baseline)
        return self.split_text_into_chunks(text)

    def _split_docx(self, file_path: str) -> List[TextChunk]:
        """处理docx文件"""
        doc = docx.Document(file_path)
        self.paragraphs = len(doc.paragraphs)
        structured_items = self._extract_docx_structure(doc)
        base_text = '\n'.join([item[0] for item in structured_items])
        self.words = len(base_text)

        return self.split_text_into_chunks(base_text)
        
    def _split_chm(self, file_path: str) -> List[TextChunk]:
        """处理chm文件"""
       
        full_text = self.extract_chm_text(file_path)

        text = '\n'.join(full_text)
        self.words = len(text)
        return self.split_text_into_chunks(text)
    
    def extract_chm_text(self, file_path: str):
        """使用7z命令行工具解压CHM"""
        temp_dir = tempfile.mkdtemp()
        try:
            # 使用7z解压CHM文件
            subprocess.run(['7z', 'x', file_path, f'-o{temp_dir}'], check=True)
            
            blocks = []
            encodings = ['utf-8', 'gbk', 'gb2312', 'big5', 'latin1']
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    if file.lower().endswith(('.htm', '.html')):
                        self.paragraphs += 1
                        file_path = os.path.join(root, file)
                        rel_path = os.path.relpath(file_path, temp_dir)

                        for encoding in encodings:
                            try:
                                with open(file_path, 'r', encoding=encoding) as f:
                                    content = f.read()
                                    blocks.append(content)
                                    break
                            except UnicodeDecodeError:
                                continue
            
            shutil.rmtree(temp_dir)
            return blocks
        except Exception as e:
            if 'temp_dir' in locals():
                shutil.rmtree(temp_dir, ignore_errors=True)
            print(f"使用7z处理CHM失败: {str(e)}")
            return []
        
    def split_text_into_chunks(self, text: str):
        if not text:
            return []
        paragraphs = text.split("\n")
        chunks = []
        current_chunk = ''
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(current_chunk) + len(para) + 1 <= self.chunk_size:
                current_chunk += (para + '\n') if current_chunk else para
            else:
                if len(para) > self.chunk_size:
                    if current_chunk:
                        chunks.append(TextChunk(
                            file_path=self.file_path,
                            text=current_chunk,
                            length=len(current_chunk),
                            language=self.language,
                            file_name=self.file_name,
                            scan_items=self.scan_items,
                            metadata={"strategy": "baseline_fixed"}
                        ))
                        current_chunk = ''
                    chunks.extend(self.split_long_paragraph(para))
                else:
                    if current_chunk:
                        chunks.append(TextChunk(
                            file_path=self.file_path,
                            text=current_chunk,
                            length=len(current_chunk),
                            language=self.language,
                            file_name=self.file_name,
                            scan_items=self.scan_items,
                            metadata={"strategy": "baseline_fixed"}
                        ))
                        current_chunk = para + '\n'
        
        if current_chunk:
            chunks.append(TextChunk(
                file_path=self.file_path,
                text=current_chunk,
                length=len(current_chunk),
                language=self.language,
                file_name=self.file_name,
                scan_items=self.scan_items,
                metadata={"strategy": "baseline_fixed"}
            ))
        return chunks

    def split_long_paragraph(self, para):
        """
            分割超长的单个段落

            Args:
                para: 要分割的长段落
                max_chunk: 每个段落的最大长度
            
            Returns:
                分割后的段落列表
        """
        if len(para) <= self.chunk_size:
            return [TextChunk(
                file_path=self.file_path,
                length=len(para),
                text=para,
                language=self.language,
                file_name=self.file_name,
                scan_items=self.scan_items,
                metadata={"strategy": "baseline_fixed"}
            )]
            
        # 尝试按句子分割
        sentences = []
        current_sentence = ""

        # 简单的句子分割逻辑(按中文句号、问号、感叹号)
        for char in para:
            current_sentence += char
            if char in ('。', '？', '！', '\n'):
                sentences.append(current_sentence)
                current_sentence = ""

        if current_sentence:
            sentences.append(current_sentence)
        
        # 如果按句子分割后仍然有句子过长，则按字数硬分割
        final_chunks = []
        current_chunk = ""
        for sentence in sentences:
            if len(sentence) > self.chunk_size:
                if current_chunk:
                    final_chunks.append(TextChunk(
                        file_path=self.file_path,
                        length=len(current_chunk),
                        text=current_chunk,
                        language=self.language,
                        file_name=self.file_name,
                        scan_items=self.scan_items,
                        metadata={"strategy": "baseline_fixed"}
                    ))
                    current_chunk = ""
                # 硬分割
                i = 0
                while i < len(sentence):
                    # 找到最后一个可分割的部分(；，,;)
                    match = re.search(pattern=r'.*([\]，；,.;])', string=sentence[i:i+self.chunk_size])
                    # 检查到达chunk_size的最大截断值
                    if match:
                        last_index = match.start(1) + i + 1
                        # 检查从i到last_index的长度是否超过chunk_size
                        if last_index - i > self.chunk_size:
                            last_index = i + self.chunk_size
                    else:
                        # 没有找到标点符号，直接截断
                        last_index = i + self.chunk_size
                    current_chunk = sentence[i:last_index]
                    i = last_index
                    final_chunks.append(TextChunk(
                        file_path=self.file_path,
                        length=len(current_chunk),
                        text=current_chunk,
                        language=self.language,
                        file_name=self.file_name,
                        scan_items=self.scan_items,
                        metadata={"strategy": "baseline_fixed"}
                    ))
            else:
                if len(current_chunk) + len(sentence) <= self.chunk_size:
                    current_chunk += sentence
                else:
                    final_chunks.append(TextChunk(
                        file_path=self.file_path,
                        length=len(current_chunk),
                        text=current_chunk,
                        language=self.language,
                        file_name=self.file_name,
                        scan_items=self.scan_items,
                        metadata={"strategy": "baseline_fixed"}
                    ))
                    current_chunk = sentence

        if current_chunk:
            final_chunks.append(TextChunk(
                file_path=self.file_path,
                length=len(current_chunk),
                text=current_chunk,
                language=self.language,
                file_name=self.file_name,
                scan_items=self.scan_items,
                metadata={"strategy": "baseline_fixed"}
            ))

        return final_chunks


    def _extract_docx_structure(self, document) -> List[tuple]:
        structured: List[tuple] = []
        for para in document.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level = self._infer_docx_heading_level(para, text)
            structured.append((text, level))

        for table in document.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    row_dict = {}
                    for idx, value in enumerate(row_data):
                        row_dict[f"column {idx + 1}"] = value
                    table_data.append(row_dict)
            if table_data:
                structured.append((json.dumps(table_data, ensure_ascii=False), 3))
        return structured

    def _infer_docx_heading_level(self, para, text: str) -> int:
        style_name = ""
        try:
            if para.style and para.style.name:
                style_name = para.style.name
        except Exception:
            style_name = ""

        if style_name:
            match = re.search(r'Heading\s*(\d+)', style_name, re.IGNORECASE)
            if match:
                try:
                    return int(match.group(1))
                except ValueError:
                    pass

        if re.match(r'^\d+(\.\d+){1,2}\s+', text):
            return 2
        if re.match(r'^\d+\s+', text):
            return 1
        if re.match(r'^[一二三四五六七八九十]+\s*[、.]\s*', text):
            return 1
        if re.match(r'^\([一二三四五六七八九十]+\)\s*', text):
            return 2
        return 99


    def _split_excel(self, file_path: str) -> List[TextChunk]:
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
        sheet_names = excel_file.sheet_names

        all_non_empty_texts = []
        for sheet_name in sheet_names:
            print(f"正在处理工作表: {sheet_name}")
            chunk_reader = pd.read_excel(excel_file, sheet_name=sheet_name, engine='openpyxl')
            for col in chunk_reader.columns:
                item = chunk_reader[col].dropna().astype(str).tolist()
                all_non_empty_texts += item
        text = '\n'.join(all_non_empty_texts)
        self.words = len(text)
        self.paragraphs = len(sheet_names)
        return self.split_text_into_chunks(text)

class PrecisionPDFProcessor:
    """
    精确换行处理的PDF处理器
    """

    # 定义句子结束标点(可根据需要扩展)
    SENTENCE_ENDINGS = {'.', '?', '!', '。', '！', '？', '"', "'", "”", "》", "】", "」", ":", ";", "…", "~"}
    ALPHABET = "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ"

    # 定义空白字符(用于判断有效前驱字符)
    WHITESPACE_CHARS = {' ', '\t', '\n', '\r', '\f', '\v'}

    def __init__(self):
        pass

    def process_pdf(self, pdf_path: str,
                    text_clean_options: Optional[Dict] = None,
                    table_clean_options: Optional[Dict] = None):
        """
        处理PDF文件，应用精确换行处理

        返回：
        {
            'text': '清洗后的文本内容',
            'tables': [
                {
                    'columns': ['column1', 'column2', ...],
                    'data': [
                        {'column1': 'value1', 'column2': 'value2', ...},
                        ...
                    ]
                },
                ...
            ]
        }
        """
        # 设置默认选项
        text_clean_options = text_clean_options or {
            'remove_special_chars': True,
            'remove_extra_whitespace': False,
            'min_line_length': 3,
            'precise_newline_handling': False
        }

        table_clean_options = table_clean_options or {
            'header_rows': 1,
            'skip_footer': 0,
            'rename_columns': True,
            'precise_newline_handling': True
        }

        result = []

        try:
            logger.info(f"开始处理PDF文件: {pdf_path}")

            # 提取和清洗文本
            raw_text = self.extract_text_from_pdf(pdf_path)
            if raw_text:
                result.append(self.clean_extracted_text(raw_text, **text_clean_options))
            logger.info(f"文本清洗完成")

            # 提取和清洗表格
            raw_tables = self.extract_tables_from_pdf(pdf_path)
            for table in raw_tables:
                processed_table = self.process_table(table, **table_clean_options)
                if processed_table:
                    result.append(json.dumps(processed_table, ensure_ascii=False))
            
            logger.info(f"表格数据清洗完成")
            return result
        
        except Exception as e:
            logger.error(f"提取PDF文本时出错: {e}")
            return result
    
    def extract_text_from_pdf(self, pdf_path: str) -> Optional[str]:
        """从PDF中提取文本内容"""
        try:
            text = self.extract_non_table_text(pdf_path)
            return text
        
        except Exception as e:
            logger.error(f"提取PDF文本时出错: {e}")
            return None
    
    def extract_non_table_text(self, pdf_path: str, output_txt_path=None):
        """
        从PDF中提取非表格文字内容

        参数：
            pdf_path(str): PDF文件路径
            output_txt_path(str, optional): 输出文本文件路径(可选)

        返回：
            str: 提取的所有非表格文本
        """
        all_text = ""

        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                # 获取当前页所有表格
                tables = page.find_tables()

                non_table_text = self.extract_text_excluding_tables_advanced(page, tables)

                # 添加到总文本中
                all_text += f"{non_table_text}\n"

            # 如果需要保存文件
            if output_txt_path:
                with open(output_txt_path, 'w', encoding='utf-8') as f:
                    f.write(all_text)

            return all_text
    
    def extract_text_excluding_tables_advanced(self, page, tables):
        """
        高级方法排除表格文本 - 基于单词位置过滤

        参数：
            page(pdfplumber.Page): 当前页面对象
            tables(list): 当前页所有表格对象列表

        返回：
            str: 提取的非表格文本
        """
        # 获取所有单词
        words = page.extract_words()

        # 如果没有表格，直接返回所有单词组成的文本
        if not tables:
            paragraphs = self.process_pdf_word_to_paragraphs(words)
            return "\n".join(paragraphs)
        

        # 获取所有表格的边界框
        table_bboxes = [table.bbox for table in tables]

        # 过滤掉表格内的单词
        filtered_words = []
        for word in words:
            word_bbox = (word['x0'], word['top'], word['x1'], word['bottom'])

            # 检查单词是否在任何表格内
            in_table = False
            for table_bbox in table_bboxes:
                if self.bbox_overlap(word_bbox, table_bbox):
                    in_table = True
                    break

            if not in_table:
                filtered_words.append(word)
        
        paragraphs = self.process_pdf_word_to_paragraphs(filtered_words)
        return "\n".join(paragraphs)


    def bbox_overlap(self, bbox1, bbox2, margin=5):
        """
        检查两个边界框是否重叠 (带边框)
        
        参数：
            bbox1(tuple): (x0, y0, x1, y1)
            bbox2(tuple): (x0, y0, x1, y1)
            margin(float): 检测重叠的边距

        返回：
            bool: 是否重叠
        """
        # 调整边界框，增加边距
        a_x0, a_y0, a_x1, a_y1 = bbox1
        b_x0, b_y0, b_x1, b_y1 = bbox2
        
        # 检查重叠
        return not (a_x1 - margin < b_x0 + margin or
                    b_x1 - margin < a_x0 + margin or
                    a_y1 - margin < b_y0 + margin or
                    b_y1 - margin < a_y0 + margin)
    
    def extract_tables_from_pdf(self, pdf_path: str) -> List[List[List[str]]]:
        """从PDF中提取表格数据"""
        try:
            tables = []
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    page_tables = page.extract_tables()
                    if page_tables:
                        tables.extend(page_tables)
            return tables
        
        except Exception as e:
            logger.error(f"提取PDF表格时出错: {e}")
            return []
    
    def clean_extracted_text(self, text: str, 
                            remove_special_chars: bool=True,
                            remove_extra_whitespace: bool=True,
                            min_line_length: int=3,
                            precise_newline_handling: bool=True) -> str:
        """
        清洗提取的PDF文本，包含精确换行处理
        """
        # 第一步：基础清洗
        lines = text.splitlines()
        cleaned_lines = []

        for line in lines:
            line = line.strip()

            if remove_extra_whitespace:
                line = re.sub(pattern=r'\s+', repl=' ', string=line)
            
            if remove_special_chars:
                line = re.sub(r'[^\w\s.,;、；：，。<>`~《》？！“”…|:!?()[]{}【】（）\-—–\'"@#$%&*=+/\\]', '', line)
            
            if len(line) >= min_line_length:
                cleaned_lines.append(line)
        
        # 第二步：精确换行处理
        if precise_newline_handling:
            cleaned_lines = self.process_newlines_precisely("\n".join(cleaned_lines))
        else:
            cleaned_text = "\n".join(cleaned_lines)
        
        # 第三步：标准化空行(保留段落分隔)
        if remove_extra_whitespace:
            cleaned_text = re.sub(pattern=fr'\n{3,}', repl='\n\n', string=cleaned_text)
        return cleaned_text.strip()

    def process_newlines_precisely(self, text: str) -> str:
        """
        精确处理换行符:
        1. 找到所有换行符位置
        2. 检查每个换行符前最后一个非空白字符
        3. 如果是句子结束标点，保留换行
        4. 否则，直接拼接(不添加任何空白字符)
        """
        if not text:
            return text
        
        # 用于存储处理后的字符
        result = []
        i = 0
        n = len(text)

        while i < n:
            if text[i] == '\n':
                # 找到换行符前最后一个非空白字符
                last_char = None
                j = i - 1
                while j >= 0:
                    if text[j] not in self.WHITESPACE_CHARS:
                        last_char = text[j]
                        break
                    j -= 1
                
                # 判断是否保留换行
                if last_char:
                    if last_char in self.SENTENCE_ENDINGS:
                        result.append('\n')
                    if last_char in self.ALPHABET:
                        result.append('△')
                # 否则不添加任何字符（相当于删除换行）
            else:
                result.append(text[i])
            i += 1

        return ''.join(result)
    
    def process_table(self, table_data: List[List[str]],
                    header_rows: int = 1,
                    skip_footer: int = 0,
                    rename_columns: bool = True,
                    precise_newline_handling: bool = True) -> Optional[Dict[str, Any]]:
        """
        处理表格数据，对每个单元格应用精确换行处理

        返回：
        {
            'columns': ['column1', 'column2', ...],
            'data': [
                {'column1': 'value1', 'column2': 'value2', ...},
                ...
            ]
        }
        """
        try:
            if not table_data or len(table_data) <= header_rows + skip_footer:
                return None
            
            # 处理表头
            if header_rows > 0:
                headers = table_data[0]
                data_rows = table_data[header_rows:len(table_data) - skip_footer]
            else:
                headers = None
                data_rows = table_data[:len(table_data) - skip_footer]
            
            # 处理列名
            if headers and not rename_columns:
                if header_rows > 1:
                    headers = [

                        self.process_tables_newlines('\n'.join(str(table_data[i][j]) for i in range(header_rows)))
                        for j in range(len(table_data[0]))]
                column_names = []
                for header in headers:
                    column_names.append(self.process_tables_newlines(header))
            else:
                column_names = [f'column {i + 1}' for i in range(len(data_rows[0]))]
            
            # 处理表格数据
            processed_data = []
            for row in data_rows:
                processed_row = {}
                for i, cell in enumerate(row):
                    if i >= len(column_names):
                        continue

                    # 处理单元格内容
                    cell_text = str(cell) if cell is not None else ''

                    # 应用精确换行处理
                    if precise_newline_handling and cell_text:

                        cell_text = self.process_tables_newlines(cell_text)
                    
                    processed_row[column_names[i]] = cell_text.strip()

                # 只添加非空行
                if any(processed_row.values()):
                    processed_data.append(processed_row)
            
            return {
                'columns': column_names,
                'data': processed_data
            }
        
        except Exception as e:
            logger.error(f"处理表格数据时出错: {e}")
            return None

    def is_paragraph_start(self, pre_line_text, cur_line_text) -> bool:
        """
        基于文本特征判断是否可能是段落开始
        """
        # 检查常见的段落起始模式
        patterns = [
            cur_line_text.startswith(('•', '-', '*', '→', '●', '步骤')),  # 项目符号
            any(cur_line_text.startswith(str(i)) for i in range(1, 10)) and len(cur_line_text) > 1 and cur_line_text[
                1] in '. )',  # 编号
            re.match(r'^(?!\d{2,}\.)\d(?:\.\d).*$', cur_line_text) not in (None, ""),
            re.match(r'^([a-z]\..*$)', cur_line_text) not in (None, ""),
            re.match(r'^(表\s*\d.*$)', cur_line_text) not in (None, ""),
            re.match(r'.*\.+\s?\d+$', pre_line_text) not in (None, "")
        ]
        return any(patterns)

    def process_pdf_word_to_paragraphs(self, words):
        """
        处理PDF读取的文字，按照段落模式
        """
        cur_line = []
        current_paragraph = []
        paragraphs = []
        prev_line_left = []
        for i, word in enumerate(words):
            cur_text = word['text']
            pre_word = words[i - 1]
            if i == 0 or abs(
                    (words[i]['bottom'] + words[i]['top']) - (words[i - 1]['bottom'] + words[i -1]['top'])) < 10:
                if i == 0:
                    prev_line_left = word['x0']
                if abs(word['x0'] - pre_word['x1']) > 0.05:
                    cur_line.append(" " + cur_text)
                else:
                    cur_line.append(cur_text)
                continue
            else:
                # 获取每个行
                line_last_word = cur_line[-1][-1]
                if line_last_word.isascii() and line_last_word.isalpha() and cur_text[0] not in ("\\"):
                    cur_line[-1] += " "
                pre_line_text = ''.join(cur_line)
                cur_line = [cur_text]
                cur_line_left = word['x0']
            # 检查缩进判断段落
            prev_word_x1 = pre_word['x1']
            cur_word_x0 = word['x0']
            indent_diff = abs(cur_line_left - prev_line_left)
            indent_len_diff = abs(cur_word_x0 - prev_word_x1)
            prev_line_left = cur_line_left
            # 判断段落边界
            is_new_paragraph = (
                indent_diff > 40 or
                indent_len_diff < 250 or 
                self.is_paragraph_start(pre_line_text, ''.join(cur_line))
            )
            current_paragraph.append(pre_line_text)
            if is_new_paragraph:
                # 保存当前段落，开始新段落
                if len(current_paragraph) > 0:
                    paragraphs.append(''.join(current_paragraph))
                current_paragraph = []
        if current_paragraph:
            paragraphs.append(''.join(current_paragraph))
        if cur_line:
            current_paragraph.append(''.join(cur_line))
        paragraphs.append(''.join(current_paragraph))
        return paragraphs
    
    def process_tables_newlines(self, text: str) -> str:
        """
        处理换行
        1. 按照换行分割单元格内容
        2. 检查行首是否是段落开头
        3. 查询换行前后是否是网址、文件目录、奇数个引号"or"
        4. 英文字符 + 空格（英文标题左侧换行）
        5. 存在问题（表格中指令换行、英文文档左侧名称在一半时换行）
        :param text:
        :return:
        """
        if not text:
            return text
        lines = text.split('\n')
        pre_line = None
        result = []
        max_line_length = max(len(line) for line in lines)
        for line in lines:
            if pre_line is None:
                result.append(line)
                pre_line = line
                continue
            # 表格间距较小（容易将英文字母分割）
            if max_line_length < 4:
                # 下一行大写开头
                if line[0].isupper():
                    result.append(' ' + line)
                else:
                    result.append(line)
                continue
            # 段首符号
            if self.is_paragraph_start(pre_line, line):
                result.append('\n' + line)
            # 上一行行末文档的句末符号+" "
            elif pre_line[-1] in (",;:"):
                result.append(' ' + line)
            # 上一行行末是否是网址 文件 引号缺失，不变
            elif self.judge_line_end(pre_line):
                result.append(line)
            # 判断是否是断裂的文件名
            elif self.judge_broken_filename(pre_line, line):
                result.append(line)
            # 上一行行末和本行行头是否是英文单词 + " "
            elif pre_line[-1].isdigit() or pre_line[-1] in ",:;" or line[0].isdigit():
                result.append(' ' + line)
            elif pre_line[-1] in self.ALPHABET and line[0] in self.ALPHABET:
                result.append('△' + line)
            else:
                result.append(line)
            pre_line = line
        return ''.join(result)
    
    def judge_line_end(self, line):
        """
        判断一行文本的行末是否为网址、文件路径，引号奇数
        """
        # 预处理：去除行末的空白字符（空格、制表符等），避免干扰匹配
        trimmed_line = line.rstrip()
        if not trimmed_line:
            return False
        
        # 网址模式：匹配 http://xxx、https://xxx、www.xxx （行末无空格）
        URL_PATTERN = r'(https?://\S*|www\.\S*)$'
        # 文件路径模式：匹配 Windows路径（C:\xxx\yyy.xxx）、Linux路径（/xxx/yyy.xxx）、相对路径（./xxx.xxx）
        FILE_PATTERN = r'([A-Za-z]:\\[^\\]*\.[^\\.]+|/[^/]*\.[^/.]+|\./[^/.]*\.[^/.]+)$'
        # 引号
        english_quotes = len(re.findall(r'["\']', line))
        # 单独匹配中文前引号“
        chinese_total_quotes = len(re.findall(r'“”', line))
        chinese_open_quotes = len(re.findall(r'“', line))
        line_end_pattern = [
            re.search(URL_PATTERN, trimmed_line),
            re.search(FILE_PATTERN, trimmed_line),
            english_quotes % 2 != 0,
            chinese_open_quotes * 2 != chinese_total_quotes
        ]
        return any(line_end_pattern)

    def judge_broken_filename(self, prev_line: str, curr_line: str) -> bool:
        """
        判断断裂的文件名（如 expert.l + nog → expert.log）
        :param prev_line: 上一行文字
        :param curr_line: 当前行文字
        :return: 修复后的内容（若修复则返回合并后的字符串，否则返回原两行）、是否修复的标记（True/False）
        """
        # 1. 定义目标日志文件的特征后缀（可根据实际需求扩展，如 .txt/.csv 等）
        LOG_SUFFIXES = {'.log', '.txt', '.csv', '.json', '.err', '.info', '.warn', '.xml', '.error'}
        # 2. 拼接上一行和当前行（中间无多余字符，匹配文件名断裂场景）
        merged_line = prev_line + curr_line
        # 3. 上一行无任何日志后缀，且拼接后有日志后缀 → 说明需要修复
        prev_has_log = any(suffix in prev_line for suffix in LOG_SUFFIXES)
        merged_has_log = any(suffix in merged_line for suffix in LOG_SUFFIXES)
        if not prev_has_log and merged_has_log:
            return True  # 需要合并
        else:
            return False  # 无需合并


class BaselineNaiveSplitter:
    """
    Baseline-Naive: 真正的外部基线
    
    特点:
    - 解析：PyPDF2 原始提取（无任何精细化处理）
    - 分块：纯滑动窗口（类似 LangChain CharacterTextSplitter with separator=""）
    
    这是实验中用作对比的最简单基线方法。
    """
    
    def __init__(self, file_path: str = "", chunk_size: int = 1000, overlap: int = 100,
                 language: str = "cn", file_name: str = "", scan_items: List = None):
        if scan_items is None:
            scan_items = []
        self.file_path = file_path
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.language = language
        self.file_name = file_name
        self.scan_items = scan_items
    
    def _extract_text_pypdf2(self, file_path: str) -> str:
        """使用 PyPDF2 原始提取 PDF 文本"""
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except Exception as e:
            logger.error(f"[Baseline-Naive] Error extracting text from {file_path}: {e}")
            return ""
    
    def _sliding_window_split(self, text: str) -> List[TextChunk]:
        """纯滑动窗口分块（无任何智能处理）"""
        chunks = []
        if not text:
            return chunks
        
        start = 0
        text_len = len(text)
        chunk_idx = 0
        
        while start < text_len:
            end = min(start + self.chunk_size, text_len)
            chunk_text = text[start:end]
            
            chunks.append(TextChunk(
                file_path=self.file_path,
                text=chunk_text,
                length=len(chunk_text),
                language=self.language,
                file_name=self.file_name,
                scan_items=self.scan_items,
                metadata={"strategy": "baseline_naive", "chunk_idx": chunk_idx}
            ))
            
            chunk_idx += 1
            if end == text_len:
                break
            start += (self.chunk_size - self.overlap)
        
        return chunks
    
    def split_document(self, file_path: str) -> List[TextChunk]:
        """
        分块文档入口
        
        支持 PDF 和纯文本文件
        """
        if file_path.lower().endswith('.pdf'):
            text = self._extract_text_pypdf2(file_path)
        else:
            # 其他文件类型直接读取
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            except Exception as e:
                logger.error(f"[Baseline-Naive] Error reading {file_path}: {e}")
                return []
        
        return self._sliding_window_split(text)